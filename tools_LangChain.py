import os
from langchain.tools import BaseTool
from urllib.error import HTTPError

import urllib
from urllib.request import urlopen
from bs4 import BeautifulSoup
from googlesearch import search

from langchain.schema import HumanMessage, SystemMessage
from langchain.chat_models import ChatOpenAI
from youtube_transcript_api import YouTubeTranscriptApi

import PyPDF2
from pptx import Presentation
from docx import Document
import copy

import torch
torch.cuda.is_available = lambda : False

from summarizer import Summarizer
import uuid

import openai
import json
import os
import re

class TextVerifier():
    def __init__(self):

        # Load your API key from an environment variable or secret management service
        openai.api_key = os.environ['OPENAI_API_KEY']
        self.system_message = """Your are a helpful assistant. The user will give you some text and a task. 
        Your goal is to report the relevant information in the text to complete the task.
        You MUST only copy the part of the text that contains relevant information.
        
        Your response MUST follow this format:
        ```
        ${correction} #your correction to include only the useful information. "!!!" if you cannot find any useful information
        ``` 
        
        This is the user request:
        {TASK_MESSAGE}    
        """

    def run(self,text):
        # Create a list to store all the messages for context
        messages = [
            {"role": "system", "content": self.system_message},
        ]

        message = f"""This is the text : ```{text}```. 
                       Does the text provide useful information for the given user requirement?"""

        condition = False
        correction = None
        # Keep repeating the following
        for _ in range(1):
            # Add each new message to the list
            messages.append({"role": "user", "content": message})

            # Request gpt-3.5-turbo for chat completion
            response = openai.ChatCompletion.create(
                model="gpt-3.5-turbo",
                messages=messages,
                stream = False,
            # "n": self.n,
            temperature=0,
                # stop=["\nDone","\tDone"]
            )

            # Print the response and add it to the messages list
            chat_message = response['choices'][0]['message']['content']
            try:
                chat_message_ = re.search(r'```(.*?)```', chat_message, re.DOTALL).group(1).strip()
            except:
                pass

            # res = json.loads(chat_message.strip())

            # response = res["response"]
            # n = len('"correction":')
            correction = chat_message_.strip()
            # correction = res["correction"]

            # if "yes" in response.lower():
            #     condition = True

            if correction.lower == "!!!":
                correction = ""

            # if "done" in response.lower():
            #     break
            # print(f"Bot: {chat_message}")
            # messages.append({"role": "assistant", "content": chat_message})

        return correction


def save_file(filename,text):
    vuid = str(uuid.uuid4())
    filename = f"{filename}_{vuid}.txt"
    folder_save = os.getcwd()+"/tmp_files"
    try:
        os.mkdir(folder_save)
    except:
        pass

    try:
        file = open(f"{folder_save}/{filename}", 'w',encoding="utf-8")
    except:
        file = open(f"{folder_save}/{filename}", 'w')

    file.write(text)
    file.close()

    # return f"{folder_save}/{filename}"
    return filename

def SummarizeText(model,text):

    # model = Summarizer()
    # result = model(text, min_length=60)
    result = model(text,num_sentences=200,min_length=10)
    # full = ''.join(result)
    return result

def get_text_and_links_from_divs(divs,page_url,max_num_divs=None,text_verifier=None):
    # texts = set([div.get_text() for div in divs])
    texts = [div.get_text().strip() for div in divs]
    links = [div.find_all('a') for div in divs]

    if max_num_divs is None:
        max_num_divs = len(texts)

    # text = "\n".join(texts)

    # # kill all script and style elements
    # for script in soup(["script", "style"]):
    #     script.extract()  # rip it out
    #
    # # get text
    # text = soup.get_text()

    def group_lines(lines):
        # lines = set(lines)
        # lines = (line.strip() for line in texts)
        # break multi-headlines into a line each
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        # drop blank lines
        text = '\n'.join(chunk for chunk in chunks if chunk)

        text = text.replace("\xa0", " ")

        return text

    def verify_text(text_verifier,text):
        correction = text_verifier.run(text)
        return correction


    texts = set(texts)
    # break into lines and remove leading and trailing space on each
    lines = []
    total_text = ""
    n = 0
    for line in texts:
        if len(line.split(" ")) > 5:
            # if len(lines):
            #     text_so_far = group_lines(lines)
            #at least n words
            lines.append(line)

            text = group_lines(lines)
            if len(lines) > max(max_num_divs/10,5) and text_verifier is not None:
                text = verify_text(text_verifier,text)
                total_text = total_text+f"\n\n {text}"
                lines[:] = []

            n = n+1
            # #Aadd link info
            # if len(link):
            #     link_url = link.get('href')
            #     if link_url is not None:
            #         if not link_url.startswith("https://") or not link_url.startswith("http://"):
            #             link_url = page_url + link_url[1:]
            #
            #     link_info = f" Additional info at '{link_url}'"
            #     lines[-1] = lines[-1]+link_info
            if n > max_num_divs:
                break

    return total_text

def parse_divs(divs,exclude=['header', 'footer', 'aside', 'script', 'style', 'nav', 'iframe', 'form','ad',"widget","button","drop"]):
    # filtered_elements = [element for element in all_elements
    #                      if element.name in excluded_tags_and_classes
    #                      or element.has_attr('class') and
    #                      any(cls in excluded_tags_and_classes for cls in element['class'])]

    divs_ = []
    divs_links = []
    for div in divs:
        # div_links = div.findAll('a')

        type_ = div.get('id')
        class_ = div.get('class')
        if class_ is not None:
            type_ = " ".join(class_)

        to_exclude = False
        if type_ is not None:
            for exclude_type in exclude:
                if exclude_type in type_.lower():
                    to_exclude = True
                    break

        if not to_exclude:
            divs_.append(div)

            # links_urls = []
            # for link in div_links:
            #     link_url = link.get('href')
            #
            #     if link_url is not None:
            #         links_urls.append(link_url)
            #
            # divs_links.append(links_urls)

    return divs_


class WebVideoTextExtractor(BaseTool):
    name = "WebVideoTextExtractor"
    description = """Takes as input a video url and returns its transcripts. It MUST only bue used on video urls.
    Input type: string #the video url
    Output type: string #the name of the temporary saved file
    """

    text_summarizer: Summarizer = Summarizer()

    def get_transcript(self, video_id):
        try:
            transcript = YouTubeTranscriptApi.get_transcript(video_id)
            return transcript
        except Exception as e:
            print(f"Error retrieving transcript: {e}")
            return None

    def _run(self, video_url):
        video_id = video_url.split("v=")[1]

        # Get the transcript
        transcript = self.get_transcript(video_id)

        text_timesteps = {'timesteps': [], 'text': []}
        if transcript is not None:
            for entry in transcript:
                text_timesteps['timesteps'].append(f"{entry['start']} - {entry['start'] + entry['duration']}")
                text_timesteps['text'].append(entry['text'])

            text = ".".join(text_timesteps['text'])

            text = SummarizeText(self.text_summarizer,text)
            # text = f"The video contains the following text:\n {text}"
        else:
            text = "Could not parse the video transcript."

        filename = save_file("video_transcript", text)

        return f"The video content is saved in the temporary file: {filename} "

    def __call__(self, *args, **kwargs):
        return self.run(*kwargs.values())


class WebPageTextExtractor(BaseTool):
    name = "WebPageTextExtractor"
    description = """Takes as input a website url and returns all the text on that website. It MUST only be used on web pages urls.
    Input type: string #the url
    Output type: string #the name of the temporary saved file
    """

    text_summarizer: Summarizer = Summarizer()
    text_verifier:TextVerifier = TextVerifier()

    def _run(self, url):
        # url = "http://news.bbc.co.uk/2/hi/health/2284783.stm"

        self.text_verifier.system_message = self.text_verifier.system_message.replace("{TASK_MESSAGE}",f"Get relevant information from the webpage at '{url}'")

        try:
            req = urllib.request.Request(url, headers={'User-Agent': "Magic Browser"})
            con = urllib.request.urlopen(req)
            html = con.read()
            # html = urlopen(url).read()
            soup = BeautifulSoup(html, features="html.parser")

            divs = soup.find_all("div")
            text = ""
            if len(divs):
                divs_ = parse_divs(divs)

                text = get_text_and_links_from_divs(divs_,url,max_num_divs=None,text_verifier=self.text_verifier)

        except HTTPError as err:
            text = f"Error reading URL. Error code: {err.code}"

        if text == "":
            text = "No available text from webpage"
        else:
            text = SummarizeText(self.text_summarizer, text)
            text = f"The webpage contains the following text:\n {text}"

        filename = save_file("webpage_content", text)
        return f"The webpage content is saved in the temporary file: {filename} "

    def __call__(self, *args, **kwargs):
        return self.run(*kwargs.values())


class FileReader(BaseTool):
    name = "FileReader"
    description = """Takes as input a user given file name (powerpoint,word, pdf) and extract the text in it.
    It MUST NOT be used with temporary .txt files generated from other tools.
    Input type: string #the filename (must not be a txt file)
    Output type: string #the name of the temporary saved file
    """

    text_summarizer: Summarizer = Summarizer()
    def _run(self, file_path):
        _, file_extension = os.path.splitext(file_path)
        file_extension = file_extension.lower()
        try:
            if file_extension == '.txt':
                with open(file_path, 'r', encoding='utf-8') as file:
                    text = file.read()

            elif file_extension == '.pdf':
                text = ''
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page_num in range(pdf_reader.numPages):
                        page = pdf_reader.getPage(page_num)
                        text += page.extractText()

            elif file_extension == '.docx':
                doc = Document(file_path)
                text = ''
                for paragraph in doc.paragraphs:
                    text += paragraph.text + '\n'

            elif file_extension == '.pptx':
                prs = Presentation(file_path)
                text = ''
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            text += shape.text + '\n'

            else:
                text = "Unsupported file format"

            text = SummarizeText(self.text_summarizer, text)
            text = f"The file contains the following text:\n {text}"

            filename = save_file("file_content", text)
        except:
            filename = file_path

        return f"The file content is saved in the temporary file: {filename} "

    def __call__(self, *args, **kwargs):
        return self.run(*kwargs.values())

class GoogleWebSearcher(BaseTool):
    name = "GoogleWebSearcher"
    description = """Takes as input query string to search on Google and returns a list of urls with associated text from the search.
    Input type: string #the search query message. It MUST be rephrased for a proper search on Google
    Output type: string #the name of the temporary saved file
    """

    N_links:int = 5
    text_verifier: TextVerifier = TextVerifier()

    def _run(self,message_search):
        # to search
        results = ""
        for n_urls,url in enumerate(search(message_search, tld="co.in", num=self.N_links, stop=self.N_links, pause=2)):
            # link_res = {}

            self.text_verifier.system_message = self.text_verifier.system_message.replace("{TASK_MESSAGE}",
                                                                                         f"Get relevant information for a web search on {message_search}")
            try:
                req = urllib.request.Request(url, headers={'User-Agent': "Magic Browser"})
                con = urllib.request.urlopen(req)
                html = con.read()
                # html_text = html.decode('utf-8')
                soup = BeautifulSoup(html, features="html.parser")

                divs = soup.find_all("div")
                page_text = ""
                if len(divs):
                    divs_ = parse_divs(divs)

                    page_text = get_text_and_links_from_divs(divs_,url,max_num_divs=20,text_verifier=self.text_verifier)

                # links_ = soup.find_all('a')
                # # links_ = soup.find_all('a', attrs={'href': re.compile("^https://")})
                #
                # links_text = []
                # for n_links,link in enumerate(links_):
                #
                #     if n_links >= self.N_links:
                #         break
                #
                #     parents = []
                #     current_sec = link
                #     div_text = "No description provided."
                #     for _ in range(5):
                #         if current_sec is not None:
                #             parent = current_sec.parent
                #             parents.append(parent)
                #             current_sec = parent
                #
                #             if parent.text is not None:
                #                 div_text = parent.text
                #                 break
                #
                #     # display the actual urls
                #     link_text = link.text
                #     link_text = link_text.strip().replace("\n", "")
                #
                #     link_url = link.get('href')
                #     if link_url is not None:
                #         # if link_url.startswith("https://") or link_url.startswith("http://"):
                #         if not link_url.startswith("https://") or not link_url.startswith("http://"):
                #             link_url = url + link_url[1:]
                #
                #     link_res[f'sub_url_{n_links}'] = {'sub_link_url':link_url,'hyper_ref_description':div_text}
                    # links_text.append(text)
            except:
                page_text = "The web page cannot be opened"
                pass

            if len(page_text):
                results = results+f"\n The webpage at {url} contains the following information:\n {page_text}  \n"

        # text = text.replace("\xa0", " ")
        text = f"""The web search has returned the following results: 
             {results}
        """

        filename = save_file("GoogleSearch_content", text)
        return f"The GoogleSearch content is saved in the temporary file: {filename} "

    def __call__(self, *args, **kwargs):
        return self.run(*kwargs.values())

class WebLinkDiscoverer(BaseTool):
    name = "WebLinkDiscoverer"
    description = """Takes as input a web link url, navigates to it and searches for other web links within the website for additional information.
        Input type: string #the url of the website to further investigate
        Output type: string #the list of urls and descriptions within the website
        """

    N_links = 10

    def _run(self, url):

        results = {}

        html = urlopen(url).read()
        # html_text = html.decode('utf-8')

        soup = BeautifulSoup(html, features="html.parser")

        links_ = soup.find_all('a')
        # links_ = soup.find_all('a', attrs={'href': re.compile("^https://")})

        links_text = []
        link_res = {}
        for n_links, link in enumerate(links_):

            if n_links >= self.N_links:
                break

            parents = []
            current_sec = link
            div_text = "No description provided."
            for _ in range(5):
                if current_sec is not None:
                    parent = current_sec.parent
                    parents.append(parent)
                    current_sec = parent

                    if parent.text is not None:
                        div_text = parent.text
                        break

            # display the actual urls
            link_text = link.text
            link_text = link_text.strip().replace("\n", "")

            link_url = link.get('href')
            if link_url is not None:
                # if link_url.startswith("https://") or link_url.startswith("http://"):
                if not link_url.startswith("https://") or not link_url.startswith("http://"):
                    link_url = url + link_url[1:]

            link_res[f'sub_url_{n_links}'] = {'sub_link_url': link_url, 'hyper_ref_description': div_text}

        results['parent_url'] = {'link_url': url, 'sub_links_info': link_res}

        text = f"""The web search has returned the following results: 
                     ```
                     {results}
                     ```
                    """
        return text

    def __call__(self, *args, **kwargs):
        return self.run(*kwargs.values())


class Chatter(BaseTool):
    name = "Chatter"
    description = """Answers the human questions.
    Input type: string #the human message
    Output type: string #the chat response
    """
    chat = ChatOpenAI(openai_api_key=os.environ['OPENAI_API_KEY'],
                      temperature=0,
                      model_name="gpt-3.5-turbo",
                      )

    sys_msg = SystemMessage(
        content="You are a helpful assistant that answers the user's questions."
    )

    messages = [sys_msg]

    def clearHistory(self):
        self.messages[:] = []
        self.messages = [self.sys_msg]

    def addHumanMessage(self, message):
        self.messages.append(HumanMessage(
            content=message)
        )

    def _run(self, message):
        if message is not None:
            self.addHumanMessage(message)

        res = self.chat(self.messages)

        self.messages.append(SystemMessage(
            content=res.content
        )
        )

        text = res.content + """ TASK DONE. Return Final Answer"""
        return text

    def __call__(self, *args, **kwargs):
        return self.run(*kwargs.values())



# class DocumentRetriever():
#     # from langchain_community.document_loaders import TextLoader
#     from langchain_core.documents import Document
#
#     # loader = TextLoader("../../modules/state_of_the_union.txt")
#     # documents = loader.load()
#     file_path = "./"
#     text = "This is the file."
#     metadata = {"source": file_path}
#     documents = [Document(page_content=text, metadata=metadata)]
#
#     from langchain_community.vectorstores import FAISS
#     from langchain_community.embeddings.openai import OpenAIEmbeddings
#     from langchain.text_splitter import CharacterTextSplitter
#     from langchain.tools.retriever import create_retriever_tool
#
#     text_splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=0)
#     texts = text_splitter.split_documents(documents)
#     embeddings = OpenAIEmbeddings()
#     db = FAISS.from_documents(texts, embeddings)
#     retriever = db.as_retriever()
#
#     retr_tool = create_retriever_tool(retriever, name="document retriever", description="Searches and returns excerpts from the 2022 State of the Union.")
#

if __name__ == "__main__":
    websearch = GoogleWebSearcher()
    #
    # webpage_extract = WebPageTextExtractor()
    # webpage_extract.run("https://www.coindesk.com/policy/2024/03/15/hong-kong-regulator-says-crypto-exchange-mexc-has-been-operating-without-a-license/")


    res = websearch.run("crypto news today")

    # BestSummarizeText(res)
